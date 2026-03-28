"""
pptx 폴더의 PPTX 파일들의 시작페이지 번호를 연속되도록 설정
- firstSlideNum 설정으로 파일 간 연속 페이지 번호 유지
"""

import os
import re
import glob
import xlsxwriter
from pathlib import Path
from datetime import datetime
from dotenv import load_dotenv
from pptx import Presentation

# 프로젝트 루트의 .env 로드 (src/ 의 상위 폴더)
_ROOT = Path(__file__).parent.parent
load_dotenv(_ROOT / ".env")

# 환경 변수
PPTX_DIR_ENV  = os.getenv("PPTX_DIR", "./pptx")
PPTX_DIR      = _ROOT / PPTX_DIR_ENV
OUTPUT_DIR    = _ROOT / os.getenv("OUTPUT_DIR", "./output")


def set_first_slide_num(prs, num):
    prs._element.set("firstSlideNum", str(num))


def count_slides(filepath):
    return len(Presentation(filepath).slides)


def get_chapter_from_filename(filename):
    """
    파일명에서 로마자 챕터(알파벳 I~X 또는 특수기호 Ⅰ~Ⅹ)를 찾아 반환합니다.
    영문 알파벳이 아닌 문자(한글, 공백, 기호, 언더바 등)와 인접한 경우만 식별합니다.
    """
    match = re.search(r'(?:^|[^a-zA-Z])(Ⅹ|Ⅸ|Ⅷ|Ⅶ|Ⅵ|Ⅴ|Ⅳ|Ⅲ|Ⅱ|Ⅰ|X|IX|VIII|VII|VI|V|IV|III|II|I)(?:[^a-zA-Z]|$)', filename)
    if match:
        return match.group(1)
    return None


def save_index_excel(rows, out_dir):
    """목차페이지정보.xlsx 생성"""
    xlsx_path = out_dir / "목차페이지정보.xlsx"
    wb = xlsxwriter.Workbook(str(xlsx_path))
    ws = wb.add_worksheet("목차")

    # 스타일
    hdr_fmt = wb.add_format({
        "bold": True, "align": "center", "valign": "vcenter",
        "bg_color": "#4472C4", "font_color": "#FFFFFF",
        "border": 1, "font_size": 11,
    })
    cell_fmt = wb.add_format({
        "align": "left", "valign": "vcenter", "border": 1, "font_size": 10,
    })
    num_fmt = wb.add_format({
        "align": "center", "valign": "vcenter", "border": 1, "font_size": 10,
    })

    # 헤더
    headers = ["파일경로", "파일명", "슬라이드 수", "시작페이지", "끝페이지"]
    col_widths = [30, 60, 12, 12, 12]
    for col, (h, w) in enumerate(zip(headers, col_widths)):
        ws.write(0, col, h, hdr_fmt)
        ws.set_column(col, col, w)
    ws.set_row(0, 20)

    # 데이터
    for row_idx, (folder, fname, n_slides, start, end) in enumerate(rows, start=1):
        ws.write(row_idx, 0, folder,   cell_fmt)
        ws.write(row_idx, 1, fname,    cell_fmt)
        ws.write(row_idx, 2, n_slides, num_fmt)
        ws.write(row_idx, 3, start,    num_fmt)
        ws.write(row_idx, 4, end,      num_fmt)
        ws.set_row(row_idx, 18)

    wb.close()
    return xlsx_path


def main():
    # 하위 폴더까지 재귀적으로 탐색 (.pptx)
    files = sorted([str(p) for p in PPTX_DIR.rglob("*.pptx")])

    if not files:
        print(f"파일 없음: {PPTX_DIR}")
        return

    # 타임스탬프 출력 폴더 생성
    timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    out_dir = OUTPUT_DIR / timestamp
    out_dir.mkdir(parents=True, exist_ok=True)

    print(f"설정값: PPTX_DIR={PPTX_DIR}")
    print(f"출력 폴더: {out_dir}")
    print(f"총 {len(files)}개 파일\n")

    print("[1/2] 슬라이드 수 파악 중...")
    slide_counts = [count_slides(f) for f in files]

    print("\n[2/2] 마스터 슬라이드 번호 설정 중...")
    page_start = 1
    current_folder = None
    current_chapter = None
    index_rows = []
    
    for i, filepath in enumerate(files):
        rel_path = Path(filepath).relative_to(PPTX_DIR)
        
        # 파일경로 표시 구조 생성 (예: ./pptx, ./pptx/III.개요)
        if str(rel_path.parent) == '.':
            folder_name = PPTX_DIR_ENV
        else:
            # 윈도우 경로(\)를 슬래시(/)로 치환해 일관성 유지
            folder_name = f"{PPTX_DIR_ENV}/{str(rel_path.parent).replace(os.sep, '/')}"
            
        fname = rel_path.name

        # === A3 예외 처리 ===
        if "A3" in fname:
            # 엑셀에만 명시하고 페이지 번호 부여 및 파일 복사(저장)는 건너뜀
            index_rows.append((folder_name, fname, slide_counts[i], "-", "-"))
            print(f"  [{folder_name}] {fname}  (A3: 페이지 번호 제외)")
            continue

        chapter = get_chapter_from_filename(fname)
        
        # 페이지 초기화 조건
        if folder_name != current_folder:
            # 폴더가 다르면 1번부터 시작
            page_start = 1
            current_folder = folder_name
            current_chapter = chapter if chapter else None
        elif not chapter:
            # 로마자 챕터가 없으면 개별 파일로 취급해 1번부터
            page_start = 1
            current_chapter = None
        elif chapter != current_chapter:
            # 폴더는 같지만 챕터가 달라지면 1번부터
            page_start = 1
            current_chapter = chapter
            
        prs = Presentation(filepath)

        set_first_slide_num(prs, page_start)

        # 서브폴더 구조 유지 (미러링)
        file_out_dir = out_dir / rel_path.parent
        file_out_dir.mkdir(parents=True, exist_ok=True)
        out_path = file_out_dir / fname
        
        prs.save(str(out_path))

        page_end = page_start + slide_counts[i] - 1
        index_rows.append((folder_name, fname, slide_counts[i], page_start, page_end))

        print(f"  [{folder_name}] {fname}")
        chapter_str = f" [Chapter {chapter}]" if chapter else ""
        print(f"    페이지 {page_start} ~ {page_end}  ({slide_counts[i]}슬라이드){chapter_str}")
        page_start = page_end + 1

    # 목차 엑셀 저장
    xlsx_path = save_index_excel(index_rows, out_dir)
    print(f"\n완료! 처리 결과가 저장되었습니다.  →  {out_dir}")
    print(f"목차 엑셀: {xlsx_path.name}")


if __name__ == "__main__":
    main()
